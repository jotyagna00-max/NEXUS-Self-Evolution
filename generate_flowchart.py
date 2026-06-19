from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
EMERALD = RGBColor(16, 185, 129)
DARK_EMERALD = RGBColor(5, 80, 55)
BLUE = RGBColor(59, 130, 246)
DARK_BG = RGBColor(10, 12, 14)
GRAY = RGBColor(156, 163, 175)
LIGHT_GRAY = RGBColor(209, 213, 219)
RED = RGBColor(239, 68, 68)
PURPLE = RGBColor(139, 92, 246)
YELLOW = RGBColor(234, 179, 8)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, shape_type, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=10, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER, font_name="JetBrains Mono"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_arrow(slide, start_left, start_top, end_left, end_top, color=EMERALD, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1, start_left, start_top, end_left, end_top
    )
    connector.line.fill.solid()
    connector.line.fill.fore_color.rgb = color
    connector.line.width = width
    return connector

def add_rounded_box(slide, left, top, width, height, text, fill=DARK_EMERALD, text_color=WHITE, font_size=9, bold=False):
    shape = add_shape(slide, left, top, width, height, MSO_SHAPE.ROUNDED_RECTANGLE, fill_color=fill, line_color=EMERALD, line_width=Pt(1))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = text_color
    run.font.bold = bold
    run.font.name = "JetBrains Mono"
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    shape.text_frame.margin_left = Pt(6)
    shape.text_frame.margin_right = Pt(6)
    shape.text_frame.margin_top = Pt(4)
    shape.text_frame.margin_bottom = Pt(4)
    return shape

# ── SLIDE 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
    "NEXUS SELF-EVOLUTION SYSTEM", 36, EMERALD, True, font_name="Orbitron")
add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.6),
    "Multi-Agent AI Orchestration Architecture", 18, GRAY, False, font_name="Space Grotesk")
add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(2),
    "A gamified self-mastery dashboard that orchestrates specialized AI agents\n(SAGE, TITAN, CHRONOS, MANAGER) to accelerate personal evolution\nthrough adaptive quests, real-time stat tracking, and neural-linked coaching.",
    12, LIGHT_GRAY, False, font_name="Space Grotesk")

# Decorative line
add_shape(slide, Inches(4), Inches(3.5), Inches(5), Pt(2), MSO_SHAPE.RECTANGLE, fill_color=EMERALD)

# ── SLIDE 2: High-Level Architecture ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
    "HIGH-LEVEL ARCHITECTURE", 22, EMERALD, True, font_name="Orbitron")

# Layer 1: Frontend
add_rounded_box(slide, Inches(4.5), Inches(0.9), Inches(4), Inches(0.6),
    "REACT FRONTEND (Vite + Tailwind v4)", DARK_EMERALD, WHITE, 10, True)

arrows_y = Inches(1.6)

# Layer 2: React Context
add_rounded_box(slide, Inches(3.5), Inches(1.8), Inches(6), Inches(0.5),
    "GameContext.tsx  (Central State - localStorage persistence)", RGBColor(20, 40, 60), BLUE, 9, True)

# Layer 3: Two AI Pipelines
add_rounded_box(slide, Inches(1), Inches(2.7), Inches(5), Inches(0.5),
    "Pipeline A: Direct Agent API\nagentService.ts → openaiAgentService.ts", RGBColor(40, 20, 60), PURPLE, 8, True)
add_rounded_box(slide, Inches(7), Inches(2.7), Inches(5.5), Inches(0.5),
    "Pipeline B: Smart Trainer Router\ntrainerService.ts → routeToAgent()", RGBColor(40, 20, 60), PURPLE, 8, True)

# Layer 4: Agents
add_rounded_box(slide, Inches(0.5), Inches(3.6), Inches(3), Inches(0.5),
    "SAGE  (Cognitive)", RGBColor(5, 50, 80), BLUE, 9, True)
add_rounded_box(slide, Inches(3.8), Inches(3.6), Inches(3), Inches(0.5),
    "TITAN  (Physical)", RGBColor(50, 15, 15), RED, 9, True)
add_rounded_box(slide, Inches(7.2), Inches(3.6), Inches(3), Inches(0.5),
    "CHRONOS  (Temporal)", RGBColor(50, 40, 10), YELLOW, 9, True)
add_rounded_box(slide, Inches(10.5), Inches(3.6), Inches(2.5), Inches(0.5),
    "MANAGER  (Orchestrator)", RGBColor(5, 55, 35), EMERALD, 9, True)

# Layer 5: NVIDIA API
add_rounded_box(slide, Inches(4), Inches(4.5), Inches(5), Inches(0.5),
    "NVIDIA Nemotron LLM (via OpenAI SDK)", RGBColor(30, 20, 40), GRAY, 10, True)

# Layer 6: Components
components_y = Inches(5.4)
comps = [
    ("InitialAssessment", Inches(0.3)), ("Dashboard", Inches(2.3)),
    ("QuestBoard", Inches(4.3)), ("TrainingHub", Inches(5.8)),
    ("PersonalizedTrainer", Inches(7.5)), ("AppControlPanel", Inches(9.8)),
    ("AgentDashboard", Inches(11.5))
]
for text, left in comps:
    add_rounded_box(slide, left, components_y, Inches(1.5), Inches(0.4),
        text, RGBColor(20, 20, 30), LIGHT_GRAY, 7, False)

# Layer 7: Storage
add_rounded_box(slide, Inches(4), Inches(6.2), Inches(5), Inches(0.5),
    "localStorage  (Zero Backend - Fully Client-Side)", RGBColor(15, 30, 20), EMERALD, 9, True)

# ── SLIDE 3: Data Flow ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
    "DATA FLOW & ALGORITHM SEQUENCE", 22, EMERALD, True, font_name="Orbitron")

# Step boxes
steps = [
    ("1", "User Action", "Click / Type / Voice\nin React UI", Inches(0.5)),
    ("2", "GameContext", "State update →\nlocalStorage sync", Inches(2.8)),
    ("3", "Route Decision", "trainerService.routeToAgent()\nkeyword-matches primary agent", Inches(5.1)),
    ("4", "Agent Execution", "SAGE | TITAN | CHRONOS | MANAGER\nbuilds system prompt", Inches(7.4)),
    ("5", "LLM Inference", "NVIDIA Nemotron\nstream/non-stream response", Inches(9.6)),
    ("6", "UI Update", "Recharts/D3 graphs\nAgent Activity Log", Inches(11.6)),
]
for num, title, desc, left in steps:
    s = add_rounded_box(slide, left, Inches(1), Inches(1.8), Inches(1.2),
        f"{num}. {title}\n{desc}", DARK_EMERALD if int(num) % 2 else RGBColor(20, 40, 60), WHITE, 8, True)

# Arrow annotations
add_text_box(slide, Inches(0.5), Inches(2.5), Inches(12), Inches(0.5),
    "┬──────────────────────────────────────────────────────────────────────────────┬", 9, EMERALD)

# Flow description
flow_steps = [
    "Initialization → App.tsx checks hasCompletedAssessment in localStorage",
    "  ├─ False → Render InitialAssessment (3-question calibration gate)",
    "  └─ True  → Render Dashboard (tabs: overview, trainer, training, quests, agents, feedback, control, settings)",
    "",
    "User interacts with AI features:",
    "  1. PersonalizedTrainer: User types message → streamTrainerRequest() → routeToAgent()",
    "     → generates coordinated prompt → streamAgentResponse() → openai.chat.completions.create()",
    "     → chunks stream back via callback → progressive UI update",
    "",
    "  2. Quest Generation: QuestGeneratorAgent builds challenges from stat thresholds + user profile",
    "     → RewardPenaltyAgent adjusts difficulty based on miss count",
    "",
    "  3. Graph Analysis: User clicks analyze on HexGraph/StatGraph → SAGE/TITAN agent interprets data",
    "     → analysis overlay with agent's textual insight",
    "",
    "  4. App Control: Lock distracting apps → visual willpower commitment → localStorage toggle",
]
for i, line in enumerate(flow_steps):
    c = EMERALD if line.startswith("  ") else (BLUE if line.startswith("  ├") or line.startswith("  └") else GRAY)
    b = True if "Initialization" in line or "User interacts" in line else False
    add_text_box(slide, Inches(0.5), Inches(2.8 + i * 0.35), Inches(12), Inches(0.35),
        line, 8, c, b, PP_ALIGN.LEFT, "JetBrains Mono")

# ── SLIDE 4: Agent System Prompts ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
    "MULTI-AGENT SYSTEM & PROMPT ARCHITECTURE", 20, EMERALD, True, font_name="Orbitron")

agents_info = [
    ("MANAGER", "Supreme Orchestrator & Life Coach", EMERALD,
     "Stoic, analytical coach. Designs hyper-personalized plans.\nDelegates to specialists. Addresses user as 'Operator'.\nHandles: motivation, emotional, multi-domain planning."),
    ("SAGE", "Intelligence & Cognitive Specialist", BLUE,
     "Deep focus on neuroplasticity, memory, speed-reading.\nAcademic tone. Handles: learning strategies, logic,\ncognitive enhancement protocols."),
    ("TITAN", "Strength & Physical Specialist", RED,
     "Sports science, hypertrophy, recovery, nutrition.\nHandles: workout plans, biological optimization,\nvitality metrics, protein timing."),
    ("CHRONOS", "Temporal & Agility Specialist", YELLOW,
     "Schedule optimization, flow-state entry, reaction speed.\nHandles: time management, habit timing,\nmobility routines, efficiency."),
]
for i, (name, role, color, desc) in enumerate(agents_info):
    y = Inches(1 + i * 1.5)
    add_rounded_box(slide, Inches(0.5), y, Inches(3), Inches(1.2),
        f"{name}\n{role}", RGBColor(10, 10, 15), color, 9, True)
    add_text_box(slide, Inches(3.8), y + Inches(0.1), Inches(9), Inches(1.2),
        desc, 9, LIGHT_GRAY, False, PP_ALIGN.LEFT, "Space Grotesk")

# ── SLIDE 5: Routing Algorithm ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
    "SMART ROUTING ALGORITHM  (trainerService.ts: routeToAgent)", 18, EMERALD, True, font_name="Orbitron")

routing = [
    ("User Input Keywords", "Primary Agent", "Secondary Agents"),
    ("workout, gym, strength, muscle,\nexercise, training, lift, cardio, physical", "TITAN", "CHRONOS"),
    ("focus, learn, study, memory,\nbrain, intelligence, read, cognitive, mental", "SAGE", "CHRONOS"),
    ("schedule, time, routine, daily,\nplan, organize, timing, flow", "CHRONOS", "SAGE, TITAN"),
    ("nutrition, diet, food, protein,\nrecovery, sleep, energy", "TITAN", "SAGE"),
    ("motivation, emotional, stress,\nconfidence, fear, anxiety, willpower, discipline", "MANAGER", "SAGE, TITAN, CHRONOS"),
    ("Default (no keyword match)", "MANAGER", "SAGE, TITAN, CHRONOS"),
]

for row_idx, (keywords, primary, secondary) in enumerate(routing):
    y = Inches(1.1 + row_idx * 0.85)
    is_header = row_idx == 0
    bg = RGBColor(5, 55, 35) if is_header else RGBColor(15, 18, 20)
    tc = WHITE if is_header else LIGHT_GRAY
    fs = 9 if is_header else 8
    b = is_header
    add_rounded_box(slide, Inches(0.5), y, Inches(5), Inches(0.7),
        keywords, bg, tc, fs, b)
    add_rounded_box(slide, Inches(5.8), y, Inches(2.5), Inches(0.7),
        primary, bg, EMERALD if not is_header else WHITE, fs, b)
    add_rounded_box(slide, Inches(8.6), y, Inches(4), Inches(0.7),
        secondary, bg, BLUE if not is_header else WHITE, fs, b)

# ── SLIDE 6: State Management ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
    "STATE MANAGEMENT & PERSISTENCE", 22, EMERALD, True, font_name="Orbitron")

add_text_box(slide, Inches(0.5), Inches(1), Inches(12), Inches(0.4),
    "All state lives in GameContext.tsx, persisted to localStorage on every change via useEffect", 10, GRAY, False, PP_ALIGN.LEFT, "Space Grotesk")

state_items = [
    ("UserStats", "strength, intelligence, agility,\nvitality, willpower, social", "InitialAssessment sets baseline\nTrainingHub protocols increment"),
    ("UserProfile", "name, goals, barriers,\nlearning style, schedule", "Set during onboarding\nUpdated via AI coach insights"),
    ("Quests / Tasks", "daily tasks + destiny quests\nwith exp/stat rewards", "QuestGeneratorAgent creates\ncompleteQuest/completeTask adjusts stats"),
    ("AppPermissions", "locked app list\n(Instagram, Twitter, TikTok)", "AppControlPanel toggles\nlocalStorage + visual commitment"),
    ("ChatHistory", "last 10 interactions\nwith agent responses", "sendCommandToAgent stores\nstreamCommandToAgent stores"),
    ("Protocols", "training protocols\n(Neural Overclock, Hypertrophy)", "Created in TrainingHub\naddProtocol function"),
]

for i, (name, content, source) in enumerate(state_items):
    y = Inches(1.6 + i * 0.9)
    add_rounded_box(slide, Inches(0.5), y, Inches(2.5), Inches(0.75),
        name, RGBColor(5, 40, 30), EMERALD, 9, True)
    add_text_box(slide, Inches(3.3), y + Inches(0.05), Inches(4), Inches(0.7),
        content, 8, LIGHT_GRAY, False, PP_ALIGN.LEFT, "JetBrains Mono")
    add_text_box(slide, Inches(7.6), y + Inches(0.05), Inches(5), Inches(0.7),
        source, 8, BLUE, False, PP_ALIGN.LEFT, "Space Grotesk")

# ── SLIDE 7: Component Tree ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
    "COMPONENT TREE & RENDERING", 22, EMERALD, True, font_name="Orbitron")

tree_lines = [
    "App.tsx",
    " └─ GameProvider (Context Wrapper)",
    "     └─ Dashboard (hasCompletedAssessment ?)",
    "         ├─ False → InitialAssessment.tsx",
    "         │          3-question calibration gate",
    "         │          Sets baseline UserStats",
    "         │",
    "         └─ True → Dashboard Layout",
    "             ├─ Header (MenuOverlay, System Status)",
    "             ├─ Main Content (tab-driven)",
    "             │   ├─ overview   → HexGraph + StatGraph + AgentActivity",
    "             │   ├─ trainer    → PersonalizedTrainer.tsx (multi-agent chat)",
    "             │   ├─ training   → TrainingHub.tsx (protocols)",
    "             │   ├─ quests     → QuestBoard.tsx (daily + destiny)",
    "             │   ├─ agents     → AgentDashboard.tsx (orchestrator status)",
    "             │   ├─ feedback   → Feedback.tsx",
    "             │   ├─ control    → AppControlPanel.tsx (app locking)",
    "             │   └─ settings   → ApiKeySettings.tsx",
    "             ├─ Floating Mic → SupremeCommander.tsx (voice overlay)",
    "             └─ localStorage sync on every state change",
]
for i, line in enumerate(tree_lines):
    c = EMERALD if "─" in line or "│" in line else (BLUE if "tsx" in line else LIGHT_GRAY)
    b = True if "App.tsx" in line or "Dashboard" in line else False
    add_text_box(slide, Inches(0.5), Inches(0.9 + i * 0.3), Inches(12), Inches(0.3),
        line, 8, c, b, PP_ALIGN.LEFT, "JetBrains Mono")

# ── SLIDE 8: Security & Key Flow ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
    "SECURITY & API KEY ARCHITECTURE", 22, EMERALD, True, font_name="Orbitron")

sec_items = [
    ("Gemini API Key", "Injected at build time via Vite define\nprocess.env.GEMINI_API_KEY\nUsed by trainerService.ts", RGBColor(5, 50, 80), BLUE),
    ("NVIDIA API Key", "Stored in localStorage\nSet via ApiKeySettings UI\nRead by openaiAgentService.ts", RGBColor(50, 15, 15), RED),
    ("Discord Webhook", "Server-side only (server.ts)\nPOST /api/community/sync\nOptional community sharing", RGBColor(30, 20, 40), PURPLE),
    ("localStorage", "All user data persisted locally\nNo backend database\nZero external data leakage", RGBColor(15, 30, 20), EMERALD),
]
for i, (name, desc, bg, color) in enumerate(sec_items):
    y = Inches(1.2 + i * 1.3)
    add_rounded_box(slide, Inches(0.5), y, Inches(3), Inches(1),
        name, bg, color, 10, True)
    add_text_box(slide, Inches(3.8), y + Inches(0.1), Inches(8.5), Inches(0.9),
        desc, 9, LIGHT_GRAY, False, PP_ALIGN.LEFT, "JetBrains Mono")

# ── SLIDE 9: Purpose & Philosophy ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
    "PURPOSE & PHILOSOPHY", 22, EMERALD, True, font_name="Orbitron")

purpose_text = [
    "CORE INTENTION",
    "Translate extreme self-discipline into a high-engagement digital workspace.",
    "No social features, no bloated gamification — pure analytical self-improvement.",
    "",
    "TARGET AUDIENCE",
    "Self-optimization enthusiasts, biohackers, deep-focus practitioners,",
    "students of strategy, and fans of calculated, narrative-driven tools.",
    "",
    "DESIGN PRINCIPLES",
    "Local-First: Your data stays on your machine. No cloud, no tracking.",
    "Client-Side AI: Multi-agent orchestration runs entirely in your browser.",
    "Anti-Telemetry: Zero verbose logs, zero metadata collection.",
    "HUD Aesthetic: Dark obsidian, emerald glow, tactical scanlines —",
    "           every pixel reinforces the 'system calibration' metaphor.",
    "",
    "ALGORITHMIC GOAL",
    "Convert daily actions → stat increments → rank progression →",
    "hunter rank ascension (E → D → C → B → A → S → SS → SSS).",
    "Each completed quest, task, or protocol directly modifies your",
    "neural evolution graph — visible proof of your self-mastery journey.",
]

for i, line in enumerate(purpose_text):
    if not line:
        continue
    c = EMERALD if line.isupper() and len(line) > 3 else LIGHT_GRAY
    b = line.isupper() and len(line) > 3
    fs = 11 if b else 9
    add_text_box(slide, Inches(0.5), Inches(1 + i * 0.35), Inches(12), Inches(0.35),
        line, fs, c, b, PP_ALIGN.LEFT, "Space Grotesk" if not b else "Orbitron")

prs.save("C:\\Users\\JOTYAGNA\\Documents\\NEXUS Self-Evolution\\NEXUS_Architecture.pptx")
print("PPT saved successfully!")
